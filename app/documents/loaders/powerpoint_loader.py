"""Carregador de apresentacoes PowerPoint (secao 5.4 do prompt mestre).

Extrai titulo, conteudo textual, listas, tabelas e notas do apresentador de
cada slide, mantendo a associacao do conteudo ao respectivo slide.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from app.core.exceptions import CorruptedFileError
from app.documents.loaders.base import DocumentLoader
from app.schemas.document import DocumentSection, ExtractedDocument, ExtractionWarning


def _table_to_text(table) -> str:
    rows_text = []
    for row_index, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        prefix = "Cabecalho" if row_index == 0 else f"Linha {row_index}"
        rows_text.append(f"{prefix}: " + " | ".join(cells))
    return "\n".join(rows_text)


class PowerPointLoader(DocumentLoader):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def load(self, file_path: Path) -> ExtractedDocument:
        try:
            presentation = Presentation(str(file_path))
        except PackageNotFoundError as exc:
            raise CorruptedFileError(
                f"Nao foi possivel abrir a apresentacao: {file_path.name}"
            ) from exc

        sections: list[DocumentSection] = []
        warnings: list[ExtractionWarning] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_had_content = False

            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title_text = title_shape.text_frame.text.strip()
                if title_text:
                    sections.append(
                        DocumentSection(text=f"Titulo: {title_text}", slide=slide_index)
                    )
                    slide_had_content = True

            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        sections.append(DocumentSection(text=text, slide=slide_index))
                        slide_had_content = True
                elif shape.has_table:
                    table_text = _table_to_text(shape.table)
                    if table_text.strip():
                        sections.append(
                            DocumentSection(
                                text=table_text,
                                slide=slide_index,
                                table_name=f"Tabela do slide {slide_index}",
                            )
                        )
                        slide_had_content = True

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    sections.append(
                        DocumentSection(
                            text=f"Notas do apresentador: {notes_text}", slide=slide_index
                        )
                    )
                    slide_had_content = True

            if not slide_had_content:
                warnings.append(
                    ExtractionWarning(
                        code="empty_slide",
                        message="Slide sem conteudo textual extraivel.",
                        location=f"slide {slide_index}",
                    )
                )

        return ExtractedDocument(
            source_path=str(file_path),
            document_format="pptx",
            sections=sections,
            warnings=warnings,
            raw_metadata={"slide_count": str(len(presentation.slides))},
        )
