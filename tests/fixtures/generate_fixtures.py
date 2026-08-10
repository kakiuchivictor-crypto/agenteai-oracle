"""Gera os documentos fictícios usados como fixtures de teste dos carregadores.

Uso (a partir da raiz do projeto, com o venv ativo):
    python tests/fixtures/generate_fixtures.py

Os arquivos gerados sao pequenos e contem apenas dados FICTICIOS (secao 33
do prompt mestre: "Nao inclua dados reais ou confidenciais"). Este script
nao precisa rodar em CI — os binarios gerados sao versionados em
`tests/fixtures/documents/`; execute-o novamente apenas se precisar
regenerar os fixtures.
"""

from __future__ import annotations

import json
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "documents"


def generate_pdf() -> None:
    import pymupdf

    doc = pymupdf.open()
    page1 = doc.new_page()
    page1.insert_text(
        (72, 72),
        "Politica de Reembolso\n\n"
        "A solicitacao de reembolso deve ser realizada em ate 7 dias corridos\n"
        "apos a compra, mediante contato com o suporte.",
    )
    page2 = doc.new_page()
    page2.insert_text(
        (72, 72),
        "Secao 2: Estornos\n\n"
        "Compras internacionais possuem prazo de analise de ate 14 dias uteis.",
    )
    doc.save(FIXTURES_DIR / "sample_policy.pdf")
    doc.close()

    # PDF sem nenhuma camada de texto (simula pagina escaneada) para testar
    # a deteccao de necessidade de OCR.
    blank_doc = pymupdf.open()
    blank_page = blank_doc.new_page()
    blank_page.draw_rect(pymupdf.Rect(50, 50, 200, 200), fill=(0, 0, 0))
    blank_doc.save(FIXTURES_DIR / "sample_scanned.pdf")
    blank_doc.close()


def generate_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Politica de Privacidade", level=1)
    doc.add_paragraph("Esta politica descreve como os dados dos clientes sao tratados.")
    doc.add_heading("Coleta de dados", level=2)
    doc.add_paragraph("Coletamos apenas os dados necessarios para a prestacao do servico.")
    doc.add_paragraph("Nome completo", style="List Bullet")
    doc.add_paragraph("Endereco de e-mail", style="List Bullet")

    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Campo"
    table.rows[0].cells[1].text = "Finalidade"
    table.rows[1].cells[0].text = "E-mail"
    table.rows[1].cells[1].text = "Comunicacao com o cliente"

    doc.save(FIXTURES_DIR / "sample_word.docx")


def generate_xlsx() -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Planos"
    sheet.append(["Plano", "Preco mensal", "Limite de usuarios"])
    sheet.append(["Profissional", "R$ 199,90", 20])
    sheet.append(["Empresarial", "R$ 499,90", 100])
    workbook.save(FIXTURES_DIR / "sample_plans.xlsx")


def generate_pptx() -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Plataforma SaaS"
    slide1.placeholders[1].text = "Visao geral dos planos disponiveis"
    slide1.notes_slide.notes_text_frame.text = "Mencionar o periodo de teste gratuito."

    slide2 = presentation.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Cancelamento"
    slide2.placeholders[1].text = "O cancelamento pode ser feito a qualquer momento"

    presentation.save(FIXTURES_DIR / "sample_deck.pptx")


def generate_markdown() -> None:
    content = """# Guia da Plataforma

Bem-vindo ao guia oficial.

## Primeiros passos

Siga as etapas abaixo:

- Crie sua conta
- Confirme o e-mail
- Acesse o painel

## Exemplo de configuracao

```
modo: producao
timeout: 30
```

## Perguntas frequentes

Consulte a secao de suporte para mais detalhes.
"""
    (FIXTURES_DIR / "sample_readme.md").write_text(content, encoding="utf-8")
    (FIXTURES_DIR / "empty.md").write_text("", encoding="utf-8")


def generate_csv() -> None:
    content = (
        "id,produto,preco,estoque\n"
        "1,Camiseta,49.90,120\n"
        "2,Calca,89.90,45\n"
        "3,Jaqueta,199.90,0\n"
    )
    (FIXTURES_DIR / "sample_data.csv").write_text(content, encoding="utf-8")


def generate_json() -> None:
    data = {
        "empresa": "Loja Exemplo",
        "planos": {
            "profissional": {"preco": 199.90, "limite_usuarios": 20},
            "empresarial": {"preco": 499.90, "limite_usuarios": 100},
        },
        "regioes_atendidas": ["Sudeste", "Sul", "Nordeste"],
        "contatos": [
            {"nome": "Suporte", "email": "suporte@example.com"},
            {"nome": "Financeiro", "email": "financeiro@example.com"},
        ],
    }
    (FIXTURES_DIR / "sample_config.json").write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (FIXTURES_DIR / "empty.json").write_text("", encoding="utf-8")


def generate_html() -> None:
    content = """<!DOCTYPE html>
<html lang="pt-br">
<head>
<meta charset="utf-8">
<title>FAQ</title>
<style>body { color: black; }</style>
<script>console.log("ignorar");</script>
</head>
<body>
<nav>Menu de navegacao repetitivo</nav>
<header>Cabecalho do site</header>
<h1>Perguntas Frequentes</h1>
<p>Respostas as duvidas mais comuns sobre nossos servicos.</p>
<h2>Pagamentos</h2>
<p>Aceitamos cartao de credito e boleto bancario.</p>
<ul>
<li>Cartao de credito</li>
<li>Boleto bancario</li>
</ul>
<h2>Planos</h2>
<table>
<tr><th>Plano</th><th>Preco</th></tr>
<tr><td>Basico</td><td>R$ 49,90</td></tr>
<tr><td>Profissional</td><td>R$ 199,90</td></tr>
</table>
<footer>Rodape repetitivo</footer>
</body>
</html>
"""
    (FIXTURES_DIR / "sample_page.html").write_text(content, encoding="utf-8")


def generate_corrupted_files() -> None:
    # Cabecalho valido de PDF (passa na verificacao de MIME/assinatura), mas
    # com estrutura interna corrompida (falha ao ser aberto pelo PyMuPDF) —
    # simula corrupcao real detectada durante a extracao, nao no upload.
    (FIXTURES_DIR / "corrupted.pdf").write_bytes(
        b"%PDF-1.4\n% estrutura interna corrompida de proposito\ntrailer garbage sem xref valido"
    )
    (FIXTURES_DIR / "corrupted.docx").write_bytes(b"isto nao e um DOCX valido")


def main() -> None:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    generate_pdf()
    generate_docx()
    generate_xlsx()
    generate_pptx()
    generate_markdown()
    generate_csv()
    generate_json()
    generate_html()
    generate_corrupted_files()
    print(f"Fixtures geradas em {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
